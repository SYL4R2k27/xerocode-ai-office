"""AI Slides — generate presentations with AI + images via Pollinations."""
from __future__ import annotations

import io
import json
import urllib.parse

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.auth import get_current_user
from app.core.database import get_db

router = APIRouter(prefix="/slides", tags=["ai-slides"])


# ---------------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------------

class SlidesGenerateRequest(BaseModel):
    topic: str = Field(..., min_length=3, max_length=2000)
    slides_count: int = Field(default=8, ge=6, le=20)
    theme: str = Field(default="business")  # business | startup | education | minimal
    language: str = Field(default="ru")  # ru | en
    include_images: bool = Field(default=True)


class SlideOutlineRequest(BaseModel):
    topic: str = Field(..., min_length=3, max_length=2000)
    slides_count: int = Field(default=8, ge=6, le=20)
    language: str = Field(default="ru")


# ---------------------------------------------------------------------------
# Theme configs
# ---------------------------------------------------------------------------

THEME_CONFIGS = {
    "business": {
        "bg_color": "1A1A2E",
        "title_color": "FFFFFF",
        "text_color": "E0E0E0",
        "accent_color": "0F3460",
    },
    "startup": {
        "bg_color": "0D1117",
        "title_color": "58A6FF",
        "text_color": "C9D1D9",
        "accent_color": "1F6FEB",
    },
    "education": {
        "bg_color": "FFFFFF",
        "title_color": "1A237E",
        "text_color": "333333",
        "accent_color": "3F51B5",
    },
    "minimal": {
        "bg_color": "FAFAFA",
        "title_color": "212121",
        "text_color": "616161",
        "accent_color": "9E9E9E",
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _hex_to_rgbstr(hex_color: str) -> str:
    """Convert hex color string to RGBColor args."""
    return hex_color


def _pollinations_url(prompt: str) -> str:
    """Build Pollinations image URL for a slide."""
    encoded = urllib.parse.quote(prompt, safe="")
    return f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=576&nologo=true"


async def _generate_slide_structure(topic: str, count: int, language: str, include_images: bool) -> list[dict]:
    """Call AI to generate slide structure as JSON."""
    from app.api.routes.documents import _call_ai, _parse_ai_json

    lang_instruction = "Пиши на русском языке." if language == "ru" else "Write in English."
    image_field = ', "image_prompt": "short English prompt for a relevant photo/illustration (max 15 words)"' if include_images else ""

    system_prompt = f"""You are an expert presentation designer. Generate a JSON array of exactly {count} slides.
Each slide object: {{"title": "slide title", "subtitle": "subtitle or empty string", "bullets": ["point 1", "point 2", "point 3"], "notes": "speaker notes"{image_field}}}
The first slide is the title slide (no bullets, just title and subtitle).
The last slide is a CTA / contacts slide.
{lang_instruction}
Be specific, use numbers and facts where possible.
Respond with ONLY a valid JSON array, no extra text."""

    content = await _call_ai(system_prompt, topic, prefer_premium=True)
    if not content:
        raise HTTPException(status_code=502, detail="AI service unavailable")

    try:
        slides = _parse_ai_json(content)
        if not isinstance(slides, list):
            raise ValueError("Expected list")
        return slides
    except (json.JSONDecodeError, ValueError) as e:
        raise HTTPException(status_code=502, detail=f"AI returned invalid JSON: {e}")


def _build_pptx(slides: list[dict], theme: str, include_images: bool) -> io.BytesIO:
    """Build PPTX file from slide data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    cfg = THEME_CONFIGS.get(theme, THEME_CONFIGS["business"])
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)

    def hex_rgb(h: str) -> RGBColor:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    for i, slide_data in enumerate(slides):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_rgb(cfg["bg_color"])

        title_text = slide_data.get("title", f"Slide {i + 1}")
        subtitle_text = slide_data.get("subtitle", "")
        bullets = slide_data.get("bullets", [])
        notes_text = slide_data.get("notes", "")
        image_prompt = slide_data.get("image_prompt", "")

        # Title
        from pptx.util import Inches, Pt
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(8.5) if (include_images and image_prompt) else Inches(11.5)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32) if i == 0 else Pt(28)
        p.font.bold = True
        p.font.color.rgb = hex_rgb(cfg["title_color"])
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        if subtitle_text:
            sub_top = Inches(1.8) if i == 0 else Inches(1.6)
            txBox2 = slide.shapes.add_textbox(left, sub_top, width, Inches(0.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(18)
            p2.font.color.rgb = hex_rgb(cfg["text_color"])
            p2.alignment = PP_ALIGN.LEFT

        # Bullets
        if bullets:
            bullet_top = Inches(2.5)
            txBox3 = slide.shapes.add_textbox(left, bullet_top, width, Inches(4.0))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            for j, bullet in enumerate(bullets):
                if j == 0:
                    p3 = tf3.paragraphs[0]
                else:
                    p3 = tf3.add_paragraph()
                p3.text = f"  {bullet}"
                p3.font.size = Pt(16)
                p3.font.color.rgb = hex_rgb(cfg["text_color"])
                p3.space_after = Pt(8)

        # Image placeholder (as a text box with URL — real image download would be async)
        if include_images and image_prompt:
            img_url = _pollinations_url(image_prompt)
            img_left = Inches(9.5)
            img_top = Inches(1.5)
            txBox4 = slide.shapes.add_textbox(img_left, img_top, Inches(3.0), Inches(0.4))
            tf4 = txBox4.text_frame
            p4 = tf4.paragraphs[0]
            p4.font.size = Pt(8)
            p4.font.color.rgb = hex_rgb(cfg["accent_color"])

            # Store image URL in slide notes for reference
            if notes_text:
                notes_text += f"\n\nImage: {img_url}"
            else:
                notes_text = f"Image: {img_url}"

        # Speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.post("/generate")
async def generate_slides(
    body: SlidesGenerateRequest,
    user=Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Generate a presentation (PPTX) with AI content and optional images."""
    slides = await _generate_slide_structure(
        topic=body.topic,
        count=body.slides_count,
        language=body.language,
        include_images=body.include_images,
    )

    # Add image URLs if requested
    if body.include_images:
        for s in slides:
            prompt = s.get("image_prompt", "")
            if prompt:
                s["image_url"] = _pollinations_url(prompt)

    buf = _build_pptx(slides, body.theme, body.include_images)

    filename = f"presentation_{body.topic[:30].replace(' ', '_')}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/outline")
async def generate_outline(
    body: SlideOutlineRequest,
    user=Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Generate slide outline only (no PPTX file)."""
    slides = await _generate_slide_structure(
        topic=body.topic,
        count=body.slides_count,
        language=body.language,
        include_images=False,
    )
    return {"slides": slides, "count": len(slides)}
