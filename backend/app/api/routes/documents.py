"""
AI Document Generation — PPTX, DOCX, XLSX, PDF
Принимает промпт → AI генерирует структуру → библиотека собирает файл.
"""
from __future__ import annotations

import io
import json
import uuid
from datetime import datetime

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.auth import get_current_user
from app.core.database import get_db
from app.models.user import User

router = APIRouter(prefix="/documents", tags=["Documents"])


# ── Schemas ─────────────────────────────────────────────────────────

class SlideRequest(BaseModel):
    prompt: str = Field(..., min_length=1, max_length=5000)
    slides_count: int = Field(default=6, ge=2, le=20)
    template: str = Field(default="business")  # business | marketing | education | minimal
    language: str = Field(default="ru")


class DocRequest(BaseModel):
    prompt: str = Field(..., min_length=1, max_length=5000)
    template: str = Field(default="report")  # report | contract | proposal | brief
    language: str = Field(default="ru")


class SpreadsheetRequest(BaseModel):
    prompt: str = Field(..., min_length=1, max_length=5000)
    language: str = Field(default="ru")


# ── Color themes ────────────────────────────────────────────────────

THEMES = {
    "business": {
        "bg": "12122A",
        "header_band": "1A1A3E",
        "title_color": "FFFFFF",
        "text_color": "B8B8D0",
        "accent": "818CF8",
        "accent2": "5EEAD4",
        "bullet_marker": "818CF8",
        "circle_accent": "1E1E48",
    },
    "marketing": {
        "bg": "0F0F14",
        "header_band": "1A1A22",
        "title_color": "FFFFFF",
        "text_color": "A8A8B8",
        "accent": "F43F5E",
        "accent2": "FBBF24",
        "bullet_marker": "F43F5E",
        "circle_accent": "2A1A1E",
    },
    "education": {
        "bg": "0C1222",
        "header_band": "142038",
        "title_color": "FFFFFF",
        "text_color": "94A3B8",
        "accent": "10B981",
        "accent2": "06B6D4",
        "bullet_marker": "10B981",
        "circle_accent": "0A2A22",
    },
    "minimal": {
        "bg": "F8F9FA",
        "header_band": "EEEEF2",
        "title_color": "111111",
        "text_color": "555555",
        "accent": "2563EB",
        "accent2": "7C3AED",
        "bullet_marker": "2563EB",
        "circle_accent": "E8E8F0",
    },
}


# ── PPTX Generation ─────────────────────────────────────────────────

def _generate_pptx_from_structure(slides_data: list[dict], theme_name: str = "business", images: dict[int, bytes] | None = None) -> io.BytesIO:
    """Build a premium PPTX with design elements: header band, accent lines, markers, watermark."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    theme = THEMES.get(theme_name, THEMES["business"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def hex_rgb(h: str) -> RGBColor:
        return RGBColor.from_string(h)

    def add_rect(slide, left, top, w, h, color, alpha=None):
        shape = slide.shapes.add_shape(1, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_rgb(color)
        shape.line.fill.background()
        return shape

    total = len(slides_data)

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        is_title_slide = idx == 0
        is_last_slide = idx == total - 1
        has_image = images and idx in images
        text_width = 7.8 if has_image else 11.5

        # ── Background ──
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_rgb(theme["bg"])

        # ── Header band ──
        band_h = Inches(1.8) if is_title_slide else Inches(1.5)
        add_rect(slide, Inches(0), Inches(0), prs.slide_width, band_h, theme["header_band"])

        # ── Accent line under header ──
        add_rect(slide, Inches(0), band_h, prs.slide_width, Pt(3), theme["accent"])

        # ── Vertical accent bar (right side) ──
        if not is_title_slide:
            bar_top = Inches(2.0)
            bar_h = Inches(3.5)
            add_rect(slide, Inches(13.1), bar_top, Pt(5), bar_h, theme["accent2"])

        # ── Title slide: decorative oval ──
        if is_title_slide:
            oval = slide.shapes.add_shape(9, Inches(9.5), Inches(-0.8), Inches(4.5), Inches(4.5))
            oval.fill.solid()
            oval.fill.fore_color.rgb = hex_rgb(theme["circle_accent"])
            oval.line.fill.background()

        # ── Title text ──
        title = slide_info.get("title", "")
        if title:
            t_top = Inches(0.4) if is_title_slide else Inches(0.3)
            t_size = Pt(40) if is_title_slide else Pt(30)
            txBox = slide.shapes.add_textbox(Inches(0.8), t_top, Inches(text_width), Inches(1.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = t_size
            p.font.bold = True
            p.font.color.rgb = hex_rgb(theme["title_color"])

        # ── Subtitle ──
        subtitle = slide_info.get("subtitle", "")
        if subtitle:
            s_top = Inches(1.15) if is_title_slide else Inches(1.0)
            txBox = slide.shapes.add_textbox(Inches(0.8), s_top, Inches(text_width), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = hex_rgb(theme["accent"])

        # ── Bullets with colored markers ──
        bullets = slide_info.get("bullets", [])
        if bullets:
            y_start = 2.0 if is_title_slide else 1.9
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(y_start), Inches(text_width - 0.5), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # Colored marker run
                run_marker = p.add_run()
                run_marker.text = "▸  "
                run_marker.font.size = Pt(18)
                run_marker.font.bold = True
                run_marker.font.color.rgb = hex_rgb(theme["bullet_marker"])
                # Text run
                run_text = p.add_run()
                run_text.text = bullet
                run_text.font.size = Pt(17)
                run_text.font.color.rgb = hex_rgb(theme["text_color"])
                p.space_after = Pt(14)

        # ── Image ──
        if has_image:
            img_stream = io.BytesIO(images[idx])
            try:
                slide.shapes.add_picture(img_stream, Inches(8.2), Inches(1.8), Inches(4.6), Inches(4.6))
            except Exception as e:
                print(f"[Documents] Image insert error slide {idx}: {e}")

        # ── Notes ──
        notes = slide_info.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        # ── Bottom bar (double layer) ──
        add_rect(slide, Inches(0), Inches(7.15), prs.slide_width, Inches(0.35), theme["header_band"])
        add_rect(slide, Inches(0), Inches(7.15), Inches(2.5), Inches(0.35), theme["accent"])

        # ── Slide number (bottom-right) ──
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.18), Inches(0.8), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{idx + 1} / {total}"
        p.font.size = Pt(9)
        p.font.color.rgb = hex_rgb(theme["accent"])
        p.alignment = PP_ALIGN.RIGHT

        # ── XeroCode watermark (bottom-left on accent bar) ──
        wm_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.18), Inches(2.0), Inches(0.3))
        tf = wm_box.text_frame
        p = tf.paragraphs[0]
        p.text = "XEROCODE AI"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = hex_rgb(theme["title_color"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


async def _call_ai(system_prompt: str, user_prompt: str, prefer_premium: bool = False) -> str | None:
    """Call AI model with fallback chain: OpenRouter (Claude/GPT) → Groq (Llama)."""
    import httpx
    from app.core.config import settings

    proxy = getattr(settings, "api_proxy", None)

    # Provider chain: direct keys first → OpenRouter → Groq fallback
    providers = []

    # 1. Anthropic Sonnet 4.7 — ПЕРВЫЙ (дешевле Opus, того же класса качество)
    if getattr(settings, "anthropic_api_key", None):
        providers.append({
            "url": "https://api.anthropic.com/v1/messages",
            "key": settings.anthropic_api_key,
            "model": "claude-sonnet-4-7",
            "name": "Anthropic/Sonnet-4.7",
            "is_anthropic": True,
            "tok": 8000,
        })

    # 2. OpenAI GPT-5.4 (прямой ключ)
    if getattr(settings, "openai_api_key", None):
        providers.append({
            "url": "https://api.openai.com/v1/chat/completions",
            "key": settings.openai_api_key,
            "model": "gpt-5.4",
            "name": "OpenAI/GPT-5.4",
            "is_openai": True,
            "tok": 4000,
        })

    # 3. OpenRouter Sonnet (fallback)
    if getattr(settings, "openrouter_api_key", None):
        providers.append({
            "url": "https://openrouter.ai/api/v1/chat/completions",
            "key": settings.openrouter_api_key,
            "model": "anthropic/claude-sonnet-4-20250514",
            "name": "OpenRouter/Sonnet",
            "tok": 4000,
        })

    # 4. Groq (free fallback)
    if getattr(settings, "groq_api_key", None):
        providers.append({
            "url": "https://api.groq.com/openai/v1/chat/completions",
            "key": settings.groq_api_key,
            "model": "llama-3.3-70b-versatile",
            "name": "Groq/Llama",
        })

    # 5. OpenRouter Gemini Flash (cheapest fallback)
    if getattr(settings, "openrouter_api_key", None):
        providers.append({
            "url": "https://openrouter.ai/api/v1/chat/completions",
            "key": settings.openrouter_api_key,
            "model": "google/gemini-2.0-flash-001",
            "name": "OpenRouter/Gemini-Flash",
        })

    for provider in providers:
        try:
            use_proxy = bool(proxy)
            transport = httpx.AsyncHTTPTransport(proxy=proxy) if use_proxy else None
            timeout = 90 if provider.get("is_anthropic") else 45
            tok = provider.get("tok", 4000)
            async with httpx.AsyncClient(transport=transport, timeout=timeout) as client:

                if provider.get("is_anthropic"):
                    resp = await client.post(
                        provider["url"],
                        headers={
                            "x-api-key": provider["key"],
                            "anthropic-version": "2023-06-01",
                            "Content-Type": "application/json",
                        },
                        json={
                            "model": provider["model"],
                            "max_tokens": tok,
                            "system": system_prompt,
                            "messages": [{"role": "user", "content": user_prompt}],
                        },
                    )
                    if resp.status_code == 200:
                        data = resp.json()
                        content = data["content"][0]["text"]
                        usage = data.get("usage", {})
                        print(f"[Documents] ✓ {provider['name']}: {usage.get('output_tokens', '?')} tokens")
                        return content
                else:
                    body: dict = {
                        "model": provider["model"],
                        "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt},
                        ],
                        "temperature": 0.7,
                    }
                    if provider.get("is_openai"):
                        body["max_completion_tokens"] = tok
                    else:
                        body["max_tokens"] = tok

                    resp = await client.post(
                        provider["url"],
                        headers={
                            "Authorization": f"Bearer {provider['key']}",
                            "Content-Type": "application/json",
                            "HTTP-Referer": "https://xerocode.ru",
                        },
                        json=body,
                    )
                    if resp.status_code == 200:
                        content = resp.json()["choices"][0]["message"]["content"]
                        print(f"[Documents] ✓ {provider['name']}")
                        return content

                print(f"[Documents] ✗ {provider['name']}: {resp.status_code} — {resp.text[:300]}")
        except Exception as e:
            print(f"[Documents] {provider['name']} error: {e}")

    return None


def _parse_ai_json(content: str) -> any:
    """Parse JSON from AI response, handling markdown fences."""
    content = content.strip()
    if content.startswith("```"):
        content = content.split("\n", 1)[1].rsplit("```", 1)[0].strip()
    return json.loads(content)


async def _ai_generate_slides(prompt: str, count: int) -> list[dict]:
    """Generate presentation slides via AI (Claude/GPT preferred)."""
    system_prompt = f"""Ты — ведущий эксперт по созданию бизнес-презентаций мирового уровня. Глубоко изучи тему и создай JSON-массив из ровно {count} слайдов.

ФОРМАТ каждого слайда:
{{"title": "заголовок (макс 8 слов)", "subtitle": "уточняющая подпись, НЕ повтор заголовка", "bullets": ["конкретный факт с цифрой или данными", "второй пункт — реальная метрика или инсайт", "третий — практический вывод"], "notes": "2-4 развёрнутых предложения для спикера: контекст, детали, что сказать аудитории", "image_prompt": "specific English prompt for high-quality photograph, max 12 words, describe scene and mood"}}

ПРАВИЛА:
- Слайд 0: ТИТУЛЬНЫЙ — только title, subtitle, image_prompt. БЕЗ bullets.
- Слайд {count - 1}: ИТОГИ/CTA — ключевые выводы и призыв к действию.
- Каждый bullet — КОНКРЕТНЫЙ ФАКТ с цифрой, процентом, названием или реальным примером. НИКАКИХ абстракций вроде «повышение эффективности» или «улучшение процессов».
- image_prompt — СПЕЦИФИЧНЫЙ для темы слайда. НЕ generic. Описывай сцену: «modern office team discussing data on screens, warm lighting» а не «business teamwork».
- Текст слайдов на РУССКОМ. image_prompt на АНГЛИЙСКОМ.
- notes: полноценные тезисы для спикера, как если бы ты готовил речь на конференции.

Отвечай ТОЛЬКО валидным JSON-массивом."""

    content = await _call_ai(system_prompt, prompt, prefer_premium=True)
    if content:
        try:
            slides = _parse_ai_json(content)
            if isinstance(slides, list) and len(slides) > 0:
                return slides[:count]
        except Exception as e:
            print(f"[Documents] JSON parse error for slides: {e}")

    return _generate_default_slides(prompt, count)


async def _ai_generate_doc(prompt: str) -> tuple[str, list[dict]]:
    """Generate document structure via AI — deep research, professional tone."""
    system_prompt = """Ты — ведущий аналитик и эксперт по созданию профессиональных документов. Глубоко изучи тему и создай JSON:

{"title": "Точное название документа", "subtitle": "Подзаголовок или дата", "sections": [{"title": "Заголовок раздела", "content": "Подробный текст раздела — минимум 3-5 предложений с конкретными данными, цифрами, примерами. Не абстракции.", "bullets": ["конкретный пункт с фактом или метрикой", "второй пункт"]}]}

ПРАВИЛА:
- Минимум 5-7 разделов. Каждый раздел — развёрнутый, с реальными данными.
- Content каждого раздела: 3-8 предложений профессионального текста с цифрами, процентами, названиями.
- Bullets: конкретные факты, НЕ абстракции. Если пишешь про рынок — укажи размер в $. Если про технологию — укажи метрики.
- Тон: профессиональный аналитический, как McKinsey или BCG отчёт.
- Текст на русском.
- Отвечай ТОЛЬКО валидным JSON."""

    content = await _call_ai(system_prompt, prompt, prefer_premium=True)
    if content:
        try:
            data = _parse_ai_json(content)
            if isinstance(data, dict) and "sections" in data:
                return data.get("title", prompt[:100]), data["sections"], data.get("subtitle", "")
        except Exception as e:
            print(f"[Documents] JSON parse error for doc: {e}")

    t, s = _generate_default_doc(prompt)
    return t, s, ""


def _generate_default_slides(prompt: str, count: int) -> list[dict]:
    """Generate slide structure without AI (fallback)."""
    slides = [
        {"title": "XeroCode AI", "subtitle": prompt[:100], "bullets": ["Сгенерировано AI", f"Промпт: {prompt[:60]}..."], "notes": "Титульный слайд"},
    ]
    # Fill remaining slides
    sections = ["Проблема", "Решение", "Как это работает", "Преимущества", "Результаты", "Следующие шаги", "Контакты"]
    for i in range(1, count):
        idx = (i - 1) % len(sections)
        slides.append({
            "title": sections[idx],
            "subtitle": "",
            "bullets": [
                f"Пункт 1 для раздела «{sections[idx]}»",
                f"Пункт 2 — детали и метрики",
                f"Пункт 3 — выводы и действия",
            ],
            "notes": f"Слайд {i + 1}: {sections[idx]}",
        })
    return slides[:count]


# ── DOCX Generation ──────────────────────────────────────────────────

def _generate_docx(title: str, sections: list[dict], subtitle: str = "") -> io.BytesIO:
    """Build a premium DOCX with professional styling."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn

    doc = Document()

    # ── Page margins ──
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # ── Title block ──
    # Accent line (table hack for colored line)
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.cell(0, 0)
    cell.text = ""
    shading = cell._element.get_or_add_tcPr()
    shading_elm = shading.makeelement(qn("w:shd"), {qn("w:fill"): "818CF8", qn("w:val"): "clear"})
    shading.append(shading_elm)
    # Make row tiny (accent line)
    row = tbl.rows[0]
    row.height = Pt(4)

    doc.add_paragraph("")

    # Title
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in heading.runs:
        run.font.size = Pt(28)
        run.font.color.rgb = RGBColor(0x18, 0x18, 0x30)
        run.font.bold = True

    # Subtitle / date
    sub_text = subtitle if subtitle else f"Подготовлено: {datetime.now().strftime('%d.%m.%Y')}"
    sub_para = doc.add_paragraph(sub_text)
    sub_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in sub_para.runs:
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x81, 0x8C, 0xF8)
        run.font.italic = True

    # Thin separator
    sep = doc.add_paragraph()
    sep.paragraph_format.space_before = Pt(6)
    sep.paragraph_format.space_after = Pt(12)
    run = sep.add_run("─" * 60)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)

    # ── Sections ──
    for i, section in enumerate(sections):
        sec_title = section.get("title", "")
        if sec_title:
            h = doc.add_heading(sec_title, level=1)
            h.paragraph_format.space_before = Pt(18)
            h.paragraph_format.space_after = Pt(6)
            for run in h.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
                run.font.bold = True

            # Accent underline for section
            accent_p = doc.add_paragraph()
            accent_p.paragraph_format.space_after = Pt(8)
            accent_run = accent_p.add_run("▬" * 12)
            accent_run.font.size = Pt(6)
            accent_run.font.color.rgb = RGBColor(0x81, 0x8C, 0xF8)

        content = section.get("content", "")
        if content:
            p = doc.add_paragraph(content)
            p.paragraph_format.space_after = Pt(8)
            p.paragraph_format.line_spacing = Pt(18)
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x44)

        bullets = section.get("bullets", [])
        for b in bullets:
            bp = doc.add_paragraph(style="List Bullet")
            bp.paragraph_format.space_after = Pt(4)
            bp.paragraph_format.left_indent = Inches(0.4)
            run = bp.add_run(b)
            run.font.size = Pt(10.5)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x55)

    # ── Footer ──
    doc.add_paragraph("")
    sep2 = doc.add_paragraph()
    run = sep2.add_run("─" * 60)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)

    footer = doc.add_paragraph()
    footer.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r1 = footer.add_run("XEROCODE AI")
    r1.font.size = Pt(8)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x81, 0x8C, 0xF8)
    r2 = footer.add_run(f"  ·  xerocode.ru  ·  {datetime.now().strftime('%d.%m.%Y %H:%M')}")
    r2.font.size = Pt(8)
    r2.font.color.rgb = RGBColor(0x99, 0x99, 0xAA)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _generate_default_doc(prompt: str) -> tuple[str, list[dict]]:
    """Default document structure."""
    return prompt[:100], [
        {"title": "Введение", "content": f"Документ создан по запросу: {prompt}", "bullets": []},
        {"title": "Основная часть", "content": "", "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"]},
        {"title": "Заключение", "content": "Выводы и рекомендации.", "bullets": []},
    ]


# ── XLSX Generation ──────────────────────────────────────────────────

def _generate_xlsx(title: str, headers: list[str], rows: list[list]) -> io.BytesIO:
    """Build a premium XLSX with professional styling."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]

    thin_border = Border(
        left=Side(style="thin", color="E0E0E8"),
        right=Side(style="thin", color="E0E0E8"),
        top=Side(style="thin", color="E0E0E8"),
        bottom=Side(style="thin", color="E0E0E8"),
    )

    # ── Title row (merged) ──
    col_count = max(len(headers), 1)
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=col_count)
    title_cell = ws.cell(row=1, column=1, value=title)
    title_cell.font = Font(bold=True, size=14, color="1A1A3E")
    title_cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 32

    # ── Subtitle row ──
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=col_count)
    sub_cell = ws.cell(row=2, column=1, value=f"Сгенерировано XeroCode AI · {datetime.now().strftime('%d.%m.%Y %H:%M')}")
    sub_cell.font = Font(size=9, italic=True, color="818CF8")
    sub_cell.alignment = Alignment(horizontal="left")
    ws.row_dimensions[2].height = 20

    # ── Empty spacer row ──
    ws.row_dimensions[3].height = 6

    # ── Header row (row 4) ──
    header_fill = PatternFill(start_color="1A1A3E", end_color="1A1A3E", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=11, name="Calibri")
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=h)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border
        # Auto-width
        ws.column_dimensions[get_column_letter(col)].width = max(14, len(str(h)) + 6)
    ws.row_dimensions[4].height = 28

    # ── Data rows ──
    even_fill = PatternFill(start_color="F4F4F8", end_color="F4F4F8", fill_type="solid")
    odd_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")
    data_font = Font(size=10.5, color="333344", name="Calibri")
    num_font = Font(size=10.5, color="1A1A3E", bold=True, name="Calibri")

    for r, row_data in enumerate(rows):
        row_num = r + 5  # data starts at row 5
        fill = even_fill if r % 2 == 0 else odd_fill
        for c, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_num, column=c, value=val)
            cell.fill = fill
            cell.border = thin_border
            cell.alignment = Alignment(vertical="center")
            if isinstance(val, (int, float)):
                cell.font = num_font
                cell.alignment = Alignment(horizontal="right", vertical="center")
                if isinstance(val, float):
                    cell.number_format = "#,##0.00"
                elif val > 999:
                    cell.number_format = "#,##0"
            else:
                cell.font = data_font
            # Update column width
            col_letter = get_column_letter(c)
            cur_width = ws.column_dimensions[col_letter].width
            val_len = len(str(val)) + 4
            if val_len > cur_width:
                ws.column_dimensions[col_letter].width = min(val_len, 40)
        ws.row_dimensions[row_num].height = 24

    # ── Bottom accent line ──
    last_row = len(rows) + 5
    for col in range(1, col_count + 1):
        cell = ws.cell(row=last_row, column=col, value="")
        cell.fill = PatternFill(start_color="818CF8", end_color="818CF8", fill_type="solid")
    ws.row_dimensions[last_row].height = 4

    # ── Freeze panes (header row) ──
    ws.freeze_panes = "A5"

    # ── Auto-filter ──
    ws.auto_filter.ref = f"A4:{get_column_letter(col_count)}{len(rows) + 4}"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ── API Endpoints ────────────────────────────────────────────────────

async def _enforce_ai_gate(db, user):
    """Central gate: must have BYOK for any provider OR platform access via plan."""
    from app.core.ai_key_resolver import resolve_key
    # Try resolve for any provider we might use; if ALL deny, reject.
    for prov in ("groq", "openrouter", "openai", "anthropic"):
        try:
            await resolve_key(db, user, prov)
            return  # at least one provider OK
        except Exception:
            continue
    from fastapi import HTTPException
    raise HTTPException(403, "Для генерации документов добавьте BYOK ключ или повысьте подписку.")


# ── Generation rate limit (daily per user) ───────────────────────────
from collections import defaultdict
from datetime import date as _date

_GEN_COUNTER: dict[str, int] = defaultdict(int)
_GEN_COUNTER_DATE: _date | None = None

# Per-plan daily limits
_PLAN_LIMITS = {
    "free": 5,
    "starter": 20,
    "pro": 100,
    "business": 500,
    "enterprise": 10_000,
}


def _daily_limit_for(user) -> int:
    plan = (getattr(user, "plan", None) or "free").lower()
    return _PLAN_LIMITS.get(plan, _PLAN_LIMITS["free"])


def _check_gen_limit(user) -> None:
    """Raise 429 if user exceeded today's document generation limit."""
    global _GEN_COUNTER_DATE
    from fastapi import HTTPException

    today = _date.today()
    if _GEN_COUNTER_DATE != today:
        _GEN_COUNTER.clear()
        _GEN_COUNTER_DATE = today

    uid = str(getattr(user, "id", "anon"))
    limit = _daily_limit_for(user)
    used = _GEN_COUNTER[uid]
    if used >= limit:
        raise HTTPException(
            status_code=429,
            detail=f"Достигнут дневной лимит генераций ({limit}). Обновите тариф или попробуйте завтра.",
        )
    _GEN_COUNTER[uid] = used + 1


@router.post("/pptx")
async def generate_pptx(
    data: SlideRequest,
    current_user: User = Depends(get_current_user),
    db=Depends(get_db),
):
    """Generate a PPTX presentation from a prompt."""
    await _enforce_ai_gate(db, current_user)
    _check_gen_limit(current_user)
    slides_data = await _ai_generate_slides(data.prompt, data.slides_count)

    # Download images from Pollinations in parallel
    import httpx
    import asyncio
    import urllib.parse
    from app.core.config import settings

    images: dict[int, bytes] = {}

    last_idx = len(slides_data) - 1

    async def _download_image(idx: int, prompt: str):
        proxy = getattr(settings, "api_proxy", None)
        transport = httpx.AsyncHTTPTransport(proxy=proxy) if proxy else None
        # Try twice: full prompt, then simplified
        for attempt in range(2):
            try:
                p = prompt if attempt == 0 else prompt.split(",")[0].strip()
                encoded = urllib.parse.quote(p, safe="")
                url = f"https://image.pollinations.ai/prompt/{encoded}?width=800&height=800&nologo=true&seed={idx}"
                async with httpx.AsyncClient(transport=transport, timeout=40) as client:
                    resp = await client.get(url)
                    if resp.status_code == 200 and len(resp.content) > 2000:
                        images[idx] = resp.content
                        print(f"[Documents] ✓ Image slide {idx}: {len(resp.content)} bytes")
                        return
                    print(f"[Documents] Image slide {idx} attempt {attempt + 1}: {resp.status_code}, {len(resp.content)} bytes")
            except Exception as e:
                print(f"[Documents] Image slide {idx} attempt {attempt + 1} error: {e}")
        print(f"[Documents] ✗ Image slide {idx}: all attempts failed")

    tasks = []
    for i, s in enumerate(slides_data):
        img_prompt = s.get("image_prompt", "")
        if img_prompt and 0 < i < last_idx:  # skip title and last CTA slide
            tasks.append(_download_image(i, img_prompt))

    if tasks:
        await asyncio.gather(*tasks)

    buf = _generate_pptx_from_structure(slides_data, data.template, images if images else None)

    filename = f"XeroCode_Slides_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/docx")
async def generate_docx(
    data: DocRequest,
    current_user: User = Depends(get_current_user),
    db=Depends(get_db),
):
    """Generate a DOCX document from a prompt."""
    await _enforce_ai_gate(db, current_user)
    _check_gen_limit(current_user)
    title, sections, subtitle = await _ai_generate_doc(data.prompt)
    buf = _generate_docx(title, sections, subtitle)

    filename = f"XeroCode_Doc_{datetime.now().strftime('%Y%m%d_%H%M')}.docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/xlsx")
async def generate_xlsx(
    data: SpreadsheetRequest,
    current_user: User = Depends(get_current_user),
    db=Depends(get_db),
):
    """Generate an XLSX spreadsheet from a prompt."""
    await _enforce_ai_gate(db, current_user)
    _check_gen_limit(current_user)

    system_prompt = """Ты — аналитик данных и эксперт по таблицам. Глубоко изучи тему и создай JSON:
{"title": "Название таблицы (макс 30 символов)", "headers": ["Колонка1", "Колонка2", ...], "rows": [["значение1", "значение2", ...], ...]}

ПРАВИЛА:
- Минимум 10-15 строк реальных данных (не заглушки).
- Числовые значения — числами (не строками): 42, 15.7, не "42".
- Данные должны быть РЕАЛИСТИЧНЫМИ и КОНКРЕТНЫМИ — настоящие названия, реальные метрики, правдоподобные цифры.
- Если тема подразумевает финансы — добавь колонки с суммами. Если аналитика — с процентами и трендами.
- Колонок 4-8, не больше 10.
- Текст на русском. Отвечай ТОЛЬКО валидным JSON."""

    headers = ["#", "Задача", "Статус", "Приоритет", "Ответственный"]
    rows = [
        [1, "Исследование рынка", "Готово", "Высокий", "AI Agent"],
        [2, "Анализ конкурентов", "В работе", "Средний", "AI Agent"],
        [3, "Подготовка отчёта", "Бэклог", "Высокий", "Оператор"],
    ]
    title = data.prompt[:31]

    content = await _call_ai(system_prompt, data.prompt, prefer_premium=True)
    if content:
        try:
            table_data = _parse_ai_json(content)
            if isinstance(table_data, dict) and "headers" in table_data and "rows" in table_data:
                headers = table_data["headers"]
                rows = table_data["rows"]
                title = table_data.get("title", title)[:31]
        except Exception as e:
            print(f"[Documents] XLSX JSON parse error: {e}")

    buf = _generate_xlsx(title, headers, rows)

    filename = f"XeroCode_Data_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/templates")
async def document_templates():
    """List available document templates."""
    return {
        "pptx": [
            {"id": "business", "name": "Бизнес", "description": "Тёмная тема, фиолетовый акцент"},
            {"id": "marketing", "name": "Маркетинг", "description": "Тёмная тема, розовый акцент"},
            {"id": "education", "name": "Образование", "description": "Тёмная тема, зелёный акцент"},
            {"id": "minimal", "name": "Минимал", "description": "Светлая тема, синий акцент"},
        ],
        "docx": [
            {"id": "report", "name": "Отчёт"},
            {"id": "contract", "name": "Договор"},
            {"id": "proposal", "name": "КП"},
            {"id": "brief", "name": "ТЗ"},
        ],
    }
